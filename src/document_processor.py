import os
import PyPDF2
from pptx import Presentation
from typing import List, Dict, Any, Optional
import logging
from sentence_transformers import SentenceTransformer
import streamlit as st

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Processes documents (PDF, PPTX) for summarization and QA.
    """
    def __init__(self, embedding_model: str = "all-MiniLM-L6-v2", 
                 chunk_size: int = 1000, chunk_overlap: int = 200,
                 process_images: bool = True,
                 use_ocr: bool = True):
        """
        Initialize the document processor.
        
        Args:
            embedding_model: Name of the sentence-transformers model to use
            chunk_size: Size of document chunks in characters (increased for better context)
            chunk_overlap: Overlap between chunks in characters (increased for better continuity)
            process_images: Whether to extract and process images
            use_ocr: Whether to use OCR for text extraction from images
        """
        self.embedding_model = SentenceTransformer(embedding_model)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.process_images = process_images
        self.use_ocr = use_ocr
        
        # Initialize image processing if enabled
        if self.process_images:
            self._setup_image_processing()
            
        # Initialize OCR if enabled
        if self.use_ocr:
            self._setup_ocr()

    def _setup_image_processing(self):
        """Set up the necessary components for image processing."""
        try:
            from transformers import VisionEncoderDecoderModel, ViTImageProcessor, AutoTokenizer
            import torch
            from PIL import Image
            import numpy as np
            
            # Check for GPU availability
            self.device = "cuda" if torch.cuda.is_available() else "cpu"
            
            # Load image captioning model (ViT + GPT-2)
            self.image_model = VisionEncoderDecoderModel.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
            self.image_model.to(self.device)
            
            # Load image processor and tokenizer
            self.image_processor = ViTImageProcessor.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
            self.tokenizer = AutoTokenizer.from_pretrained("nlpconnect/vit-gpt2-image-captioning")
            
            self.image_processing_available = True
            logger.info("Image processing successfully initialized")
        except ImportError as e:
            logger.warning(f"Image processing libraries not available: {e}")
            logger.warning("Images will be skipped during document processing")
            self.image_processing_available = False
    
    def _setup_ocr(self):
        """Set up OCR components."""
        try:
            # Check if pytesseract is available
            import pytesseract
            from PIL import Image
            
            # Test if tesseract is installed and configured
            pytesseract.get_tesseract_version()
            
            self.ocr_available = True
            logger.info("OCR processing successfully initialized")
        except (ImportError, Exception) as e:
            logger.warning(f"OCR libraries not available or Tesseract not installed: {e}")
            logger.warning("OCR will be skipped during document processing")
            self.ocr_available = False
            
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a document file and return chunks with embeddings.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of document chunks with content, embeddings, and metadata
        """
        # Extract text based on file type
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            page_texts = self._extract_pdf_text(file_path)
            # Extract topics from the combined text of all pages
            full_text = ' '.join([page_info['text'] for page_info in page_texts])
            topics = self._extract_topics(full_text)
            document_structure = self.extract_document_structure(full_text)
            
            # Process each page and track page numbers
            processed_chunks = []
            chunk_id = 0
            
            for page_info in page_texts:
                page_text = page_info['text']
                page_num = page_info['page_number']
                
                # Chunk the page text
                chunks = self._chunk_text(page_text)
                
                # Create embeddings for each chunk from this page
                for chunk in chunks:
                    embedding = self.embedding_model.encode(chunk, show_progress_bar=False)
                    
                    processed_chunks.append({
                        'content': chunk,
                        'embedding': embedding,
                        'metadata': {
                            'source': os.path.basename(file_path),
                            'chunk_id': chunk_id,
                            'topics': topics,
                            'page_number': page_num,
                            'document_structure': document_structure['title'] if chunk_id == 0 else ''
                        }
                    })
                    chunk_id += 1
            
            return processed_chunks
            
        elif file_extension in ['.pptx', '.ppt']:
            slide_texts = self._extract_pptx_text(file_path)
            # Extract topics from the combined text of all slides
            full_text = ' '.join([slide_info['text'] for slide_info in slide_texts])
            topics = self._extract_topics(full_text)
            document_structure = self.extract_document_structure(full_text)
            
            # Process each slide and track slide numbers
            processed_chunks = []
            chunk_id = 0
            
            for slide_info in slide_texts:
                slide_text = slide_info['text']
                slide_num = slide_info['slide_number']
                
                # Chunk the slide text
                chunks = self._chunk_text(slide_text)
                
                # Create embeddings for each chunk from this slide
                for chunk in chunks:
                    embedding = self.embedding_model.encode(chunk, show_progress_bar=False)
                    
                    processed_chunks.append({
                        'content': chunk,
                        'embedding': embedding,
                        'metadata': {
                            'source': os.path.basename(file_path),
                            'chunk_id': chunk_id,
                            'topics': topics,
                            'page_number': slide_num,  # For slides, use slide number as page number
                            'document_structure': document_structure['title'] if chunk_id == 0 else ''
                        }
                    })
                    chunk_id += 1
            
            return processed_chunks
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    def _extract_pdf_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text and images from PDF file with page tracking.
        
        Returns:
            List of dictionaries with text content, images, and page number
        """
        page_texts = []
        
        try:
            # Use PyMuPDF (fitz) for better PDF processing including images
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc):
                # Extract text
                text = page.get_text()
                page_info = {
                    'text': text,
                    'page_number': page_num + 1,
                    'images': []
                }
                
                # Extract images if processing is enabled and available
                if self.process_images:
                    image_list = page.get_images(full=True)
                    
                    for img_index, img_info in enumerate(image_list):
                        # Extract image
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # Process image to extract content
                        try:
                            image_result = self._process_image_bytes(image_bytes)
                            
                            # Add image information
                            page_info['images'].append({
                                'index': img_index,
                                'caption': image_result['caption'],
                                'ocr_text': image_result['ocr_text']
                            })
                            
                            # Append image information to the text with markers
                            if image_result['ocr_text']:
                                page_info['text'] += f"\n[IMAGE {img_index + 1} TEXT: {image_result['ocr_text']}]\n"
                            if image_result['caption']:
                                page_info['text'] += f"\n[IMAGE {img_index + 1} CAPTION: {image_result['caption']}]\n"
                            
                        except Exception as e:
                            logger.warning(f"Failed to process image on page {page_num + 1}: {e}")
                
                page_texts.append(page_info)
                
            return page_texts
            
        except ImportError:
            # Fall back to PyPDF2 if PyMuPDF is not available (text only)
            logger.warning("PyMuPDF not available, falling back to PyPDF2 (text-only extraction)")
            page_texts = []
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text.strip():  # Only add non-empty pages
                        page_texts.append({
                            'text': text,
                            'page_number': page_num + 1,
                            'images': []  # No images in the fallback method
                        })
            return page_texts
    
    def _extract_pptx_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text and images from PowerPoint file with slide tracking.
        
        Returns:
            List of dictionaries with text content, images, and slide number
        """
        slide_texts = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_info = {
                'text': "",
                'slide_number': slide_num + 1,
                'images': []
            }
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_info['text'] += shape.text + "\n"
                
                # Extract images if processing is enabled
                if self.process_images:
                    try:
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            # Get image data
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Process image
                            image_result = self._process_image_bytes(image_bytes)
                            
                            # Add image information
                            img_index = len(slide_info['images'])
                            slide_info['images'].append({
                                'index': img_index,
                                'caption': image_result['caption'],
                                'ocr_text': image_result['ocr_text']
                            })
                            
                            # Append image information to the text with markers
                            if image_result['ocr_text']:
                                slide_info['text'] += f"\n[IMAGE {img_index + 1} TEXT: {image_result['ocr_text']}]\n"
                            if image_result['caption']:
                                slide_info['text'] += f"\n[IMAGE {img_index + 1} CAPTION: {image_result['caption']}]\n"
                            
                    except Exception as e:
                        logger.warning(f"Failed to process image on slide {slide_num + 1}: {e}")
            
            # Only add slides with content
            if slide_info['text'].strip() or slide_info['images']:
                slide_texts.append(slide_info)
        
        return slide_texts

    def _process_image_bytes(self, image_bytes: bytes) -> Dict[str, str]:
        """
        Process image bytes to extract content using OCR and/or generate caption.
        
        Args:
            image_bytes: Raw image data
            
        Returns:
            Dictionary with OCR text and/or caption
        """
        from io import BytesIO
        from PIL import Image
        
        result = {
            'ocr_text': '',
            'caption': ''
        }
        
        try:
            # Convert bytes to PIL Image
            image = Image.open(BytesIO(image_bytes)).convert('RGB')
            
            # Perform OCR if enabled and available
            if self.use_ocr and self.ocr_available:
                try:
                    import pytesseract
                    # Extract text using OCR
                    ocr_text = pytesseract.image_to_string(image)
                    result['ocr_text'] = ocr_text.strip()
                except Exception as e:
                    logger.error(f"OCR processing error: {e}")
            
            # Generate caption if image processing is available
            if self.image_processing_available:
                try:
                    import torch
                    # Process image for the captioning model
                    pixel_values = self.image_processor(image, return_tensors="pt").pixel_values.to(self.device)
                    
                    # Generate caption
                    with torch.no_grad():
                        output_ids = self.image_model.generate(
                            pixel_values,
                            max_length=50,
                            num_beams=4,
                            early_stopping=True
                        )
                    
                    # Decode caption
                    caption = self.tokenizer.decode(output_ids[0], skip_special_tokens=True)
                    result['caption'] = caption
                except Exception as e:
                    logger.error(f"Image captioning error: {e}")
            
            return result
            
        except Exception as e:
            logger.error(f"Image processing error: {e}")
            return result
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks with overlap."""
        chunks = []
        start = 0
        
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            
            # Adjust end to avoid splitting words
            if end < len(text):
                # Look for the last space within the chunk
                last_space = text.rfind(' ', start, end)
                if last_space != -1 and last_space > start:
                   end = last_space + 1  # Include the space

            # Add the chunk
            chunks.append(text[start:end])
            
            # Move the start position for the next chunk, considering overlap
            start = max(end - self.chunk_overlap, start + 1)
        
        return chunks
    
    def extract_document_structure(self, text: str) -> Dict[str, Any]:
        """
        Extract document structure like headings, sections, and key elements.
        Useful for structured summarization.
        
        Args:
            text: Document text
            
        Returns:
            Dictionary with document structure information
        """
        import re
        
        # Identify potential headings (capitalized lines, numbered sections, etc.)
        lines = text.split('\n')
        
        structure = {
            "title": "",
            "headings": [],
            "sections": {}
        }
        
        # Try to identify document title (usually at the beginning)
        for i, line in enumerate(lines[:10]):  # Check first 10 lines
            line = line.strip()
            if line and len(line) < 100 and (line.isupper() or line[0].isupper()):
                structure["title"] = line
                break
        
        # Identify headings and sections
        current_heading = "Introduction"
        current_section = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check if this looks like a heading
            is_heading = False
            
            # Numbered heading pattern (e.g., "1. Introduction", "1.2 Methods")
            if re.match(r'^[\d\.]+\s+[A-Z]', line):
                is_heading = True
            # All caps heading
            elif line.isupper() and 3 < len(line) < 100:
                is_heading = True
            # Title case heading that's not a full sentence
            elif line[0].isupper() and "." not in line and len(line) < 100:
                is_heading = True
                
            if is_heading:
                # Save previous section
                if current_section:
                    structure["sections"][current_heading] = '\n'.join(current_section)
                    current_section = []
                    
                # Set new heading
                current_heading = line
                structure["headings"].append(current_heading)
            else:
                current_section.append(line)
        
        # Save the last section
        if current_section:
            structure["sections"][current_heading] = '\n'.join(current_section)
        
        return structure

    def _extract_topics(self, text: str, max_topics: int = 5) -> List[str]:
        """
        Extract key topics from text using K-means clustering on word embeddings.
        
        Args:
            text: The text to extract topics from
            max_topics: Maximum number of topics/clusters to extract
            
        Returns:
            List of topic keywords/phrases representing cluster centers
        """
        try:
            # Import necessary libraries
            from sklearn.cluster import KMeans
            from sklearn.feature_extraction.text import CountVectorizer
            import nltk
            import numpy as np
            
            # Download NLTK resources if not already downloaded
            try:
                nltk.data.find('tokenizers/punkt')
                nltk.data.find('corpora/stopwords')
            except LookupError:
                nltk.download('punkt', quiet=True)
                nltk.download('stopwords', quiet=True)
            
            # Get stop words
            from nltk.corpus import stopwords
            stop_words = set(stopwords.words('english'))
            
            # Tokenize and clean text
            from nltk.tokenize import word_tokenize
            
            # Split text into sentences for better context
            sentences = nltk.sent_tokenize(text)
            
            # Filter out very short sentences
            sentences = [s for s in sentences if len(s.split()) > 3]
            
            if not sentences:
                return self._simple_topic_extraction(text, max_topics)
            
            # Generate sentence embeddings
            sentence_embeddings = self.embedding_model.encode(sentences, show_progress_bar=False)
            
            # Determine appropriate number of clusters
            # If we have very few sentences, reduce number of clusters
            n_clusters = min(max_topics, len(sentences) // 2)
            n_clusters = max(1, n_clusters)  # Ensure at least one cluster
            
            # Apply K-means clustering
            kmeans = KMeans(n_clusters=n_clusters, random_state=42)
            kmeans.fit(sentence_embeddings)
            
            # Identify central sentences for each cluster (closest to centroid)
            closest_indices = []
            for i in range(n_clusters):
                # Get sentences in this cluster
                cluster_indices = np.where(kmeans.labels_ == i)[0]
                
                if len(cluster_indices) == 0:
                    continue
                    
                # Find the sentence closest to the cluster centroid
                distances = np.linalg.norm(
                    sentence_embeddings[cluster_indices] - kmeans.cluster_centers_[i].reshape(1, -1), 
                    axis=1
                )
                closest_idx = cluster_indices[np.argmin(distances)]
                closest_indices.append(closest_idx)
            
            # Extract keywords from the central sentences
            central_sentences = [sentences[i] for i in closest_indices]
            
            # Use KeyBERT if available for more meaningful keywords
            try:
                from keybert import KeyBERT
                
                # Initialize the KeyBERT model
                kw_model = KeyBERT(model=self.embedding_model)
                
                # Extract keywords from each central sentence
                all_keywords = []
                for sentence in central_sentences:
                    keywords = kw_model.extract_keywords(
                        sentence,
                        keyphrase_ngram_range=(1, 2),
                        stop_words='english',
                        use_mmr=True,
                        diversity=0.7,
                        top_n=1  # Just get the top keyword per central sentence
                    )
                    if keywords:
                        all_keywords.append(keywords[0][0])  # Add just the keyword, not the score
                
                # If we couldn't extract keywords from all clusters, fall back to extraction method
                if len(all_keywords) < n_clusters // 2:
                    return self._extract_keywords_from_sentences(central_sentences, max_topics)
                    
                return all_keywords[:max_topics]  # Return at most max_topics keywords
                
            except ImportError:
                # Fall back to simpler keyword extraction method
                return self._extract_keywords_from_sentences(central_sentences, max_topics)
                
        except (ImportError, Exception) as e:
            # Fall back to the simple method if clustering fails
            import logging
            logging.warning(f"Clustering-based topic extraction failed with error: {e}")
            return self._simple_topic_extraction(text, max_topics)

    def _extract_keywords_from_sentences(self, sentences: List[str], max_topics: int = 5) -> List[str]:
        """
        Extract keywords from a list of central sentences using TF-IDF.
        
        Args:
            sentences: List of central sentences from clusters
            max_topics: Maximum number of topics to extract
            
        Returns:
            List of topic keywords/phrases
        """
        from sklearn.feature_extraction.text import TfidfVectorizer
        import numpy as np
        
        # Combine sentences into a single corpus, but keep track of the original sentences
        corpus = sentences
        
        # Create a TF-IDF vectorizer
        vectorizer = TfidfVectorizer(
            ngram_range=(1, 2),  # Consider single words and bigrams
            stop_words='english',
            max_features=100
        )
        
        # Generate TF-IDF matrix
        try:
            tfidf_matrix = vectorizer.fit_transform(corpus)
        except ValueError:
            # Handle case when there's not enough data
            return self._simple_topic_extraction(" ".join(sentences), max_topics)
        
        # Get feature names (words or bigrams)
        feature_names = vectorizer.get_feature_names_out()
        
        # For each sentence, find the most important feature
        keywords = []
        for i, sent in enumerate(corpus):
            if i >= len(tfidf_matrix.toarray()):
                continue
                
            # Get the TF-IDF scores for this document
            tfidf_scores = tfidf_matrix.toarray()[i]
            
            # Find the top scoring features for this document
            top_indices = np.argsort(tfidf_scores)[-2:][::-1]  # Get top 2 features
            
            for idx in top_indices:
                if tfidf_scores[idx] > 0:  # Only consider non-zero scores
                    keywords.append(feature_names[idx])
        
        # Remove duplicates and limit to max_topics
        unique_keywords = list(dict.fromkeys(keywords))
        return unique_keywords[:max_topics]
        
    def _simple_topic_extraction(self, text: str, max_topics: int = 5) -> List[str]:
        """
        Fallback method for topic extraction using TF-IDF-like approach.
        
        Args:
            text: The text to extract topics from
            max_topics: Maximum number of topics to extract
            
        Returns:
            List of topic keywords
        """
        import re
        from collections import Counter
        
        # Define stopwords - a more comprehensive list than the original
        stopwords = set([
            'a', 'about', 'above', 'after', 'again', 'against', 'all', 'am', 'an', 'and', 'any', 'are', 
            'as', 'at', 'be', 'because', 'been', 'before', 'being', 'below', 'between', 'both', 'but', 
            'by', 'can', 'did', 'do', 'does', 'doing', 'don', 'down', 'during', 'each', 'few', 'for', 
            'from', 'further', 'had', 'has', 'have', 'having', 'he', 'her', 'here', 'hers', 'herself', 
            'him', 'himself', 'his', 'how', 'i', 'if', 'in', 'into', 'is', 'it', 'its', 'itself', 'just',
            'me', 'more', 'most', 'my', 'myself', 'no', 'nor', 'not', 'now', 'of', 'off', 'on', 'once', 
            'only', 'or', 'other', 'our', 'ours', 'ourselves', 'out', 'over', 'own', 's', 'same', 'she', 
            'should', 'so', 'some', 'such', 't', 'than', 'that', 'the', 'their', 'theirs', 'them', 
            'themselves', 'then', 'there', 'these', 'they', 'this', 'those', 'through', 'to', 'too', 
            'under', 'until', 'up', 'very', 'was', 'we', 'were', 'what', 'when', 'where', 'which', 
            'while', 'who', 'whom', 'why', 'will', 'with', 'you', 'your', 'yours', 'yourself', 'yourselves'
        ])
        
        # Clean text
        text = text.lower()
        text = re.sub(r'[^\w\s]', '', text)  # Remove punctuation
        
        # Tokenize
        words = text.split()
        
        # Filter out stopwords and short words
        filtered_words = [word for word in words if word not in stopwords and len(word) > 3]
        
        # Count word frequencies
        word_counts = Counter(filtered_words)
        
        # Calculate a simple TF-IDF-like score
        # For TF: normalize by document length
        doc_length = len(filtered_words)
        word_tf = {word: count/doc_length for word, count in word_counts.items()}
        
        # Since we don't have a corpus for IDF, we'll use a simple heuristic:
        # - Penalize very frequent words (might be domain-specific stopwords)
        # - Boost mid-frequency words (likely to be topical)
        max_freq = max(word_tf.values()) if word_tf else 0
        word_scores = {word: freq * (0.5 + 0.5 * (freq / max_freq)) for word, freq in word_tf.items()}
        
        # Get the top N topics
        top_topics = [word for word, _ in sorted(word_scores.items(), key=lambda x: x[1], reverse=True)[:max_topics]]
        
        return top_topics